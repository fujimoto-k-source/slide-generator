"""
Slide Generator v2 — AI-Powered Presentation Engine
Gemini Content Analysis + Auto Charts + Smart Layout + AI Imagery
"""

import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import io, re, json, math
from datetime import datetime

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as ticker
import numpy as np
from PIL import Image, ImageDraw, ImageFilter

try:
    from google import genai
    from google.genai import types as genai_types
    GEMINI_AVAILABLE = True
except ImportError:
    GEMINI_AVAILABLE = False

# =============================================
#  Page Config & CSS
# =============================================
st.set_page_config(page_title="Slide Generator", layout="centered",
                   initial_sidebar_state="collapsed")

st.markdown("""
<style>
#MainMenu, footer, header { visibility: hidden; }
.stApp { background:#FFF; font-family:'Helvetica Neue',Arial,sans-serif; }
.main .block-container { max-width:740px; padding-top:56px; padding-bottom:88px; }

.pg-eyebrow { font-size:10px; font-weight:700; letter-spacing:.2em;
              text-transform:uppercase; color:#CCC; margin-bottom:8px; }
.pg-title   { font-size:26px; font-weight:700; letter-spacing:-.3px;
              color:#111; margin-bottom:6px; line-height:1.2; }
.pg-sub     { font-size:13px; color:#999; margin-bottom:52px; line-height:1.7; }
.step-label { font-size:10px; font-weight:700; letter-spacing:.2em;
              text-transform:uppercase; color:#CCC; margin:40px 0 16px; display:block; }

div[data-testid="stRadio"] label { display:none!important; }
div[data-testid="stRadio"] > div { display:flex; flex-direction:column; gap:8px; }
div[data-testid="stRadio"] > div > label {
    border:1px solid #EAEAEA!important; border-radius:2px!important;
    padding:18px 22px!important; background:#FAFAFA!important;
    cursor:pointer!important; transition:all .15s!important; }
div[data-testid="stRadio"] > div > label:hover {
    border-color:#111!important; background:#FFF!important; }
div[data-testid="stRadio"] > div > label[data-checked="true"] {
    border-color:#111!important; background:#FFF!important;
    box-shadow:inset 3px 0 0 #111!important; }

label, .stTextArea label, .stTextInput label {
    font-size:10px!important; font-weight:700!important;
    letter-spacing:.15em!important; text-transform:uppercase!important; color:#BBB!important; }
.stTextArea textarea {
    background:#F9F9F9!important; border:1px solid #E8E8E8!important;
    border-radius:2px!important; font-size:13.5px!important; color:#111!important;
    font-family:'Helvetica Neue',Arial,sans-serif!important;
    line-height:1.7!important; padding:18px!important; resize:vertical!important; }
.stTextArea textarea:focus { border-color:#111!important; box-shadow:none!important; }
.stTextInput input {
    background:#F9F9F9!important; border:1px solid #E8E8E8!important;
    border-radius:2px!important; font-size:13.5px!important; color:#111!important;
    padding:11px 15px!important; }
.stTextInput input:focus { border-color:#111!important; box-shadow:none!important; }

.stButton>button {
    background:#111!important; color:#FFF!important; border:none!important;
    border-radius:2px!important; padding:15px 32px!important;
    font-size:11px!important; font-weight:700!important; letter-spacing:.14em!important;
    text-transform:uppercase!important; width:100%!important; transition:opacity .15s!important; }
.stButton>button:hover { opacity:.65!important; }
.stDownloadButton>button {
    background:#FFF!important; color:#111!important;
    border:1.5px solid #111!important; border-radius:2px!important;
    padding:15px 32px!important; font-size:11px!important; font-weight:700!important;
    letter-spacing:.14em!important; text-transform:uppercase!important;
    width:100%!important; margin-top:10px!important; transition:background .15s!important; }
.stDownloadButton>button:hover { background:#F5F5F5!important; }

.swatch { display:inline-block; width:13px; height:13px; border-radius:1px;
          margin-right:8px; vertical-align:middle; border:1px solid rgba(0,0,0,.08); }
.divider { border:none; border-top:1px solid #F0F0F0; margin:44px 0; }
[data-testid="column"] { padding:0 6px!important; }

.preview-card { background:#FAFAFA; border:1px solid #EAEAEA; border-radius:2px;
                padding:12px 16px; margin:6px 0; }
.preview-tag  { font-size:9px; font-weight:700; letter-spacing:.15em;
                text-transform:uppercase; color:#BBB; }
.preview-name { font-size:13px; font-weight:600; color:#222; margin:3px 0 0; }
</style>
""", unsafe_allow_html=True)

# =============================================
#  DTP Core Utilities
# =============================================
W, H = 13.33, 7.5

def _rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def _hext(h):
    h = h.lstrip('#')
    return (int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def _hexm(h):
    return h if h.startswith('#') else f'#{h}'

def _run(run, font='Calibri', size=None, bold=False, italic=False,
         color=None, tracking=0, ea='Meiryo UI'):
    run.font.name = font
    if size: run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color: run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    if tracking:
        rPr.set('spc', str(int(tracking)))
    for ch in list(rPr):
        if ch.tag == qn('a:ea'): rPr.remove(ch)
    etree.SubElement(rPr, qn('a:ea')).set('typeface', ea)

def _para(p, align=PP_ALIGN.LEFT, sb=0, sa=0, lh=1.0):
    p.alignment = align
    p.space_before = Pt(sb)
    p.space_after = Pt(sa)
    if lh != 1.0:
        pPr = p._p.get_or_add_pPr()
        for ch in list(pPr):
            if ch.tag == qn('a:lnSpc'): pPr.remove(ch)
        spc = etree.SubElement(etree.SubElement(pPr, qn('a:lnSpc')), qn('a:spcPct'))
        spc.set('val', str(int(lh * 100000)))

def _rect(sl, l, t, w, h, fill=None):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.fill.background()
    return s

def _tb(sl, txt, l, t, w, h, size=16, bold=False, italic=False,
        color=None, align=PP_ALIGN.LEFT, tracking=0, lh=1.0,
        font='Calibri', ea='Meiryo UI', sb=0):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    _para(p, align=align, sb=sb, lh=lh)
    r = p.add_run(); r.text = str(txt)
    _run(r, font=font, size=size, bold=bold, italic=italic,
         color=color, tracking=tracking, ea=ea)
    return tb

def _img(sl, data, l, t, w, h):
    if data is None: return
    if isinstance(data, Image.Image):
        buf = io.BytesIO()
        data.convert('RGB').save(buf, format='PNG')
        buf.seek(0)
        sl.shapes.add_picture(buf, Inches(l), Inches(t), Inches(w), Inches(h))
    elif isinstance(data, (bytes, bytearray)):
        sl.shapes.add_picture(io.BytesIO(data), Inches(l), Inches(t), Inches(w), Inches(h))
    elif isinstance(data, io.BytesIO):
        data.seek(0)
        sl.shapes.add_picture(data, Inches(l), Inches(t), Inches(w), Inches(h))

# =============================================
#  Theme
# =============================================
LM = 0.75
RM = 0.75
CW = W - LM - RM

ACCENT_PRESETS = {
    'Navy':     '#0A2463',
    'Gold':     '#A67C3A',
    'Crimson':  '#8C1C2E',
    'Emerald':  '#0B5E4A',
    'Slate':    '#2E3F4F',
}

def _theme(ahex):
    a = _rgb(ahex)
    return dict(
        accent=a, ahex=ahex,
        white=_rgb('FFFFFF'), ink=_rgb('0D0D0D'),
        body=_rgb('2C2C2C'), sub=_rgb('717171'),
        rules=_rgb('EBEBEB'), numc=_rgb('CCCCCC'),
        muted=_rgb('9A9A9A'), bg=_rgb('F5F5F5'),
    )

# =============================================
#  Gemini API
# =============================================
ANALYSIS_PROMPT = """You are a presentation design expert.
Analyze the following text and output a professional slide deck structure in JSON.

RULES:
1. Assign the best layout type to each slide
2. Extract chart_data when numerical data exists
3. Add image_prompt (in English, minimal corporate style) for visual slides
4. First slide MUST be "title_hero"
5. Last slide MUST be "closing"
6. NEVER omit any information from the text. Reflect ALL content.
7. All image_prompt must include: "minimal geometric corporate illustration, clean white background, subtle shadows, no text, no labels"

LAYOUT TYPES:
- title_hero: Title slide (cover)
- stat_callout: Big KPI/number emphasis (stat_value + stat_label required)
- split_image: Text left + image right (image_prompt required)
- process_flow: Steps/flow (steps array required, each with label/title/desc)
- comparison: Compare/contrast (left/right required, each with title + items array)
- chart_full: Chart-dominant (chart_data required)
- grid_cards: Multiple item cards (cards array required, each with title + desc)
- closing: Ending slide

chart_data structure:
{"type": "bar"|"line"|"pie"|"horizontal_bar", "labels": [...], "values": [...], "unit": "..."}
Multi-series: {"type":"line","labels":[...],"datasets":[{"name":"...","values":[...]}],"unit":"..."}

Output: {"slides": [...]}

TEXT:
"""

def _init_gemini(key):
    if not GEMINI_AVAILABLE:
        return None
    try:
        return genai.Client(api_key=key)
    except Exception as e:
        st.error(f'Gemini Init Error: {e}')
        return None

def _analyze(client, text):
    try:
        r = client.models.generate_content(
            model='gemini-2.0-flash',
            contents=ANALYSIS_PROMPT + text,
            config=genai_types.GenerateContentConfig(
                response_mime_type='application/json',
            ),
        )
        data = json.loads(r.text)
        return data.get('slides', data if isinstance(data, list) else [])
    except Exception as e:
        st.error(f'Analysis Error: {e}')
        return None

def _gen_image(client, prompt):
    if not client or not prompt:
        return None
    full = (
        f"Create a minimal, professional illustration for a presentation slide. "
        f"Style: clean geometric corporate art, white background, subtle gradients, "
        f"elegant shadows, absolutely no text or labels in the image. "
        f"Subject: {prompt}"
    )
    # Try Gemini image generation
    try:
        r = client.models.generate_content(
            model='gemini-2.0-flash-exp',
            contents=full,
            config=genai_types.GenerateContentConfig(
                response_modalities=['IMAGE', 'TEXT'],
            ),
        )
        for part in r.candidates[0].content.parts:
            if hasattr(part, 'inline_data') and part.inline_data and part.inline_data.data:
                return Image.open(io.BytesIO(part.inline_data.data))
    except Exception:
        pass
    # Fallback: Imagen
    try:
        r = client.models.generate_images(
            model='imagen-3.0-generate-002',
            prompt=full,
            config=genai_types.GenerateImagesConfig(number_of_images=1),
        )
        if r.generated_images:
            return r.generated_images[0].image
    except Exception:
        pass
    return None

# =============================================
#  Abstract Visual Fallback (Pillow)
# =============================================
def _abstract(ahex, w=1200, h=675):
    img = Image.new('RGBA', (w, h), (255,255,255,255))
    base = _hext(ahex)

    ov = Image.new('RGBA', (w,h), (0,0,0,0))
    od = ImageDraw.Draw(ov)

    cx, cy, r = int(w*0.72), int(h*0.62), int(w*0.38)
    for i in range(r, 0, -3):
        a = int(22 * (i/r))
        od.ellipse([cx-i, cy-i, cx+i, cy+i], fill=(*base, a))
    cx2, cy2, r2 = int(w*0.18), int(h*0.22), int(w*0.18)
    for i in range(r2, 0, -3):
        a = int(14 * (i/r2))
        od.ellipse([cx2-i, cy2-i, cx2+i, cy2+i], fill=(*base, a))

    img = Image.alpha_composite(img, ov).convert('RGB')
    img = img.filter(ImageFilter.GaussianBlur(radius=35))

    buf = io.BytesIO()
    img.save(buf, format='PNG', quality=95)
    buf.seek(0)
    return buf.getvalue()

# =============================================
#  Chart Engine (matplotlib)
# =============================================
def _sfloats(lst):
    out = []
    for v in (lst or []):
        try: out.append(float(v))
        except: out.append(0)
    return out

def _cbase(ahex, w=7, h=4):
    fig, ax = plt.subplots(figsize=(w, h), dpi=200)
    fig.patch.set_facecolor('white')
    fig.patch.set_alpha(0)
    ax.set_facecolor('white')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#E8E8E8')
    ax.spines['bottom'].set_color('#E8E8E8')
    ax.tick_params(colors='#888888', labelsize=9)
    return fig, ax

def _csave(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', bbox_inches='tight', transparent=True, dpi=200)
    plt.close(fig)
    buf.seek(0)
    return buf.getvalue()

def _cbar(data, ahex):
    labels = data.get('labels', [])
    values = _sfloats(data.get('values', []))
    if not labels or not values: return None
    fig, ax = _cbase(ahex)
    accent = _hexm(ahex)
    n = len(values)
    colors = ['#DCDCDC'] * n
    if values:
        colors[values.index(max(values))] = accent
    ax.bar(labels, values, color=colors, width=0.55, edgecolor='none')
    mx = max(values) if values else 1
    for i, (l, v) in enumerate(zip(labels, values)):
        ax.text(i, v + mx*0.02, f'{v:,.0f}', ha='center', va='bottom',
                fontsize=9, fontweight='bold', color='#333')
    ax.set_ylim(0, mx*1.18)
    u = data.get('unit', '')
    if u: ax.set_ylabel(u, fontsize=9, color='#AAA')
    fig.tight_layout(pad=0.8)
    return _csave(fig)

def _cline(data, ahex):
    labels = data.get('labels', [])
    datasets = data.get('datasets', [])
    if not datasets:
        v = _sfloats(data.get('values', []))
        if v: datasets = [{'name': '', 'values': v}]
    if not labels or not datasets: return None
    fig, ax = _cbase(ahex)
    accent = _hexm(ahex)
    palette = [accent, '#AAAAAA', '#CCCCCC', '#666666']
    for i, ds in enumerate(datasets):
        vals = _sfloats(ds.get('values', []))
        c = palette[i % len(palette)]
        ax.plot(labels[:len(vals)], vals, color=c, linewidth=2.5,
                marker='o', markersize=5, markerfacecolor='white',
                markeredgecolor=c, markeredgewidth=2)
        if vals:
            ax.text(len(vals)-1, vals[-1], f'  {vals[-1]:,.0f}',
                    va='center', fontsize=8, fontweight='bold', color=c)
    u = data.get('unit', '')
    if u: ax.set_ylabel(u, fontsize=9, color='#AAA')
    fig.tight_layout(pad=0.8)
    return _csave(fig)

def _cpie(data, ahex):
    labels = data.get('labels', [])
    values = _sfloats(data.get('values', []))
    if not labels or not values: return None
    fig, ax = _cbase(ahex, w=5, h=4)
    base = _hext(ahex)
    colors = []
    for i in range(len(values)):
        f = 1.0 - i*0.14
        c = tuple(min(255, int(v*f + 255*(1-f)*0.3)) for v in base)
        colors.append(f'#{c[0]:02x}{c[1]:02x}{c[2]:02x}')
    wedges, _, autotexts = ax.pie(
        values, labels=None, autopct='%1.0f%%', colors=colors, startangle=90,
        wedgeprops=dict(width=0.42, edgecolor='white', linewidth=2.5), pctdistance=0.78)
    for t in autotexts:
        t.set_fontsize(9); t.set_fontweight('bold'); t.set_color('#333')
    ax.legend(labels, loc='center', fontsize=8, frameon=False)
    fig.tight_layout(pad=0.5)
    return _csave(fig)

def _chbar(data, ahex):
    labels = data.get('labels', [])
    values = _sfloats(data.get('values', []))
    if not labels or not values: return None
    fig, ax = _cbase(ahex)
    accent = _hexm(ahex)
    y = range(len(labels))
    ax.barh(list(y), values, color=accent, height=0.5, edgecolor='none', alpha=0.85)
    ax.set_yticks(list(y)); ax.set_yticklabels(labels, fontsize=9)
    ax.invert_yaxis()
    mx = max(values) if values else 1
    for bar, val in zip(ax.patches, values):
        ax.text(bar.get_width() + mx*0.02, bar.get_y()+bar.get_height()/2,
                f'{val:,.0f}', va='center', fontsize=9, fontweight='bold', color='#333')
    ax.spines['left'].set_visible(False)
    fig.tight_layout(pad=0.8)
    return _csave(fig)

def _make_chart(cd, ahex):
    if not cd: return None
    t = cd.get('type', 'bar')
    fn = {'bar': _cbar, 'line': _cline, 'pie': _cpie,
          'donut': _cpie, 'horizontal_bar': _chbar}.get(t, _cbar)
    try:
        return fn(cd, ahex)
    except Exception:
        return None

# =============================================
#  Layout: Title Hero
# =============================================
def _lay_title(sl, sd, th, ahex, img=None):
    _rect(sl, 0, 0, W, 0.04, fill=th['accent'])

    if img:
        _img(sl, img, W*0.48, 0.04, W*0.52, H-0.08)
        # Fade overlay (white gradient)
        ov = Image.new('RGBA', (400, 675), (255,255,255,0))
        d = ImageDraw.Draw(ov)
        for x in range(400):
            a = int(255 * (1 - x/400))
            d.line([(x,0),(x,675)], fill=(255,255,255,a))
        buf = io.BytesIO(); ov.save(buf, format='PNG'); buf.seek(0)
        _img(sl, buf.getvalue(), W*0.42, 0.04, W*0.16, H-0.08)

    _tb(sl, sd.get('title',''), LM, 1.6, 6.8, 2.4,
        size=52, bold=True, color=th['ink'], tracking=-100, lh=1.08)
    _rect(sl, LM, 4.15, 2.0, 0.025, fill=th['accent'])
    sub = sd.get('subtitle','')
    if sub:
        _tb(sl, sub, LM, 4.40, 6.2, 0.9,
            size=16, color=th['muted'], lh=1.55)
    _rect(sl, 0, H-0.04, W, 0.04, fill=th['accent'])

# =============================================
#  Layout: Stat Callout
# =============================================
def _lay_stat(sl, sd, th, ahex, chart=None):
    _tb(sl, sd.get('title',''), LM, 0.35, CW, 0.5,
        size=11, bold=True, color=th['muted'], tracking=120)
    _rect(sl, LM, 0.92, CW, 0.008, fill=th['rules'])

    _tb(sl, sd.get('stat_value',''), LM, 1.15, 6.0, 2.6,
        size=96, bold=True, color=th['accent'], tracking=-80, lh=1.0)
    _tb(sl, sd.get('stat_label',''), LM, 3.55, 5.5, 0.5,
        size=16, color=th['sub'])

    bullets = sd.get('bullets', [])
    y = 4.25
    for b in bullets[:3]:
        _tb(sl, f'\u2014  {b}', LM, y, 5.5, 0.4,
            size=13, color=th['body'], lh=1.4)
        y += 0.48

    if chart:
        _img(sl, chart, 7.0, 1.2, 5.6, 4.6)
    _rect(sl, 0, H-0.04, W, 0.04, fill=th['accent'])

# =============================================
#  Layout: Split Image
# =============================================
def _lay_split(sl, sd, th, ahex, img=None):
    _tb(sl, sd.get('title',''), LM, 0.35, 5.8, 0.7,
        size=28, bold=True, color=th['ink'], tracking=-40, lh=1.12)
    _rect(sl, LM, 1.15, 1.5, 0.02, fill=th['accent'])

    bullets = sd.get('bullets', [])
    y = 1.50
    for b in bullets:
        _tb(sl, f'\u2014  {b}', LM, y, 5.6, 0.52,
            size=15, color=th['body'], lh=1.45)
        y += 0.62

    if img:
        _img(sl, img, 7.2, 0.45, 5.4, 6.6)
    else:
        _rect(sl, 7.2, 0.45, 5.4, 6.6, fill=th['bg'])
    _rect(sl, 0, H-0.04, W, 0.04, fill=th['accent'])

# =============================================
#  Layout: Process Flow
# =============================================
def _lay_flow(sl, sd, th, ahex):
    _tb(sl, sd.get('title',''), LM, 0.35, CW, 0.7,
        size=28, bold=True, color=th['ink'], tracking=-40, lh=1.12)
    _rect(sl, LM, 1.15, CW, 0.008, fill=th['rules'])
    _rect(sl, LM, 1.15, 1.5, 0.008, fill=th['accent'])

    steps = sd.get('steps', [])
    n = len(steps)
    if n == 0: return
    _rect(sl, 0, H-0.04, W, 0.04, fill=th['accent'])

    gap = 0.35
    sw = min(2.8, (CW - gap*(n-1)) / n)
    total = n*sw + (n-1)*gap
    sx = LM + (CW - total) / 2
    circ = 0.65
    cy_top = 2.0

    for i, step in enumerate(steps):
        x = sx + i*(sw + gap)
        cx = x + sw/2

        oval = sl.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(cx - circ/2), Inches(cy_top),
            Inches(circ), Inches(circ))
        oval.fill.solid(); oval.fill.fore_color.rgb = th['accent']
        oval.line.fill.background()

        lbl = step.get('label', f'{i+1:02d}')
        _tb(sl, lbl, cx-circ/2, cy_top+0.12, circ, circ-0.12,
            size=18, bold=True, color=th['white'], align=PP_ALIGN.CENTER)

        if i < n-1:
            ax1 = cx + circ/2 + 0.06
            ncx = sx + (i+1)*(sw+gap) + sw/2
            ax2 = ncx - circ/2 - 0.06
            _rect(sl, ax1, cy_top + circ/2 - 0.006, ax2-ax1, 0.012, fill=th['rules'])

        _tb(sl, step.get('title',''), x, cy_top + circ + 0.20, sw, 0.5,
            size=14, bold=True, color=th['ink'], align=PP_ALIGN.CENTER, tracking=-10, lh=1.12)

        desc = step.get('desc','')
        if desc:
            _tb(sl, desc, x, cy_top + circ + 0.72, sw, 1.2,
                size=11, color=th['sub'], align=PP_ALIGN.CENTER, lh=1.42)

# =============================================
#  Layout: Comparison
# =============================================
def _lay_compare(sl, sd, th, ahex):
    _tb(sl, sd.get('title',''), LM, 0.35, CW, 0.7,
        size=28, bold=True, color=th['ink'], tracking=-40, lh=1.12)
    _rect(sl, LM, 1.15, CW, 0.008, fill=th['rules'])

    hw = (CW - 0.5) / 2
    left = sd.get('left', {})
    right = sd.get('right', {})

    _rect(sl, LM, 1.45, hw, 5.3, fill=th['bg'])
    _tb(sl, left.get('title','Before'), LM+0.3, 1.72, hw-0.6, 0.5,
        size=18, bold=True, color=th['accent'], tracking=-20)
    y = 2.35
    for item in left.get('items', []):
        _tb(sl, f'\u2014  {item}', LM+0.3, y, hw-0.6, 0.4,
            size=13, color=th['body'], lh=1.4)
        y += 0.50

    mid = LM + hw + 0.18
    _rect(sl, mid, 1.65, 0.012, 4.8, fill=th['accent'])

    rx = LM + hw + 0.5
    _rect(sl, rx, 1.45, hw, 5.3, fill=th['bg'])
    _tb(sl, right.get('title','After'), rx+0.3, 1.72, hw-0.6, 0.5,
        size=18, bold=True, color=th['accent'], tracking=-20)
    y = 2.35
    for item in right.get('items', []):
        _tb(sl, f'\u2014  {item}', rx+0.3, y, hw-0.6, 0.4,
            size=13, color=th['body'], lh=1.4)
        y += 0.50

    _rect(sl, 0, H-0.04, W, 0.04, fill=th['accent'])

# =============================================
#  Layout: Chart Full
# =============================================
def _lay_chart(sl, sd, th, ahex, chart=None):
    _tb(sl, sd.get('title',''), LM, 0.35, CW-1, 0.7,
        size=28, bold=True, color=th['ink'], tracking=-40, lh=1.12)
    _rect(sl, LM, 1.15, CW, 0.008, fill=th['rules'])
    _rect(sl, LM, 1.15, 1.5, 0.008, fill=th['accent'])

    if chart:
        _img(sl, chart, LM+0.2, 1.50, CW-0.4, 5.0)

    bullets = sd.get('bullets', [])
    if bullets and not chart:
        y = 1.60
        for b in bullets:
            _tb(sl, f'\u2014  {b}', LM, y, CW, 0.4,
                size=14, color=th['body'], lh=1.45)
            y += 0.55
    _rect(sl, 0, H-0.04, W, 0.04, fill=th['accent'])

# =============================================
#  Layout: Grid Cards
# =============================================
def _lay_grid(sl, sd, th, ahex):
    _tb(sl, sd.get('title',''), LM, 0.35, CW, 0.7,
        size=28, bold=True, color=th['ink'], tracking=-40, lh=1.12)
    _rect(sl, LM, 1.15, CW, 0.008, fill=th['rules'])
    _rect(sl, LM, 1.15, 1.5, 0.008, fill=th['accent'])

    cards = sd.get('cards', [])[:6]
    n = len(cards)
    if n == 0: return
    _rect(sl, 0, H-0.04, W, 0.04, fill=th['accent'])

    cols = min(n, 3)
    rows = math.ceil(n / cols)
    gap = 0.22
    cw = (CW - gap*(cols-1)) / cols
    ch = 2.15 if rows == 1 else 1.85
    sy = 1.50

    for i, card in enumerate(cards):
        col = i % cols
        row = i // cols
        x = LM + col*(cw + gap)
        y = sy + row*(ch + gap)

        _rect(sl, x, y, cw, ch, fill=th['bg'])
        _rect(sl, x, y, cw, 0.032, fill=th['accent'])

        _tb(sl, f'{i+1:02d}', x+0.2, y+0.22, 0.5, 0.3,
            size=9, bold=True, color=th['numc'], tracking=80)
        _tb(sl, card.get('title',''), x+0.2, y+0.52, cw-0.4, 0.4,
            size=14, bold=True, color=th['ink'], tracking=-10, lh=1.15)
        desc = card.get('desc','')
        if desc:
            _tb(sl, desc, x+0.2, y+0.95, cw-0.4, ch-1.15,
                size=11, color=th['sub'], lh=1.4)

# =============================================
#  Layout: Closing
# =============================================
def _lay_close(sl, sd, th, ahex, author=''):
    _rect(sl, 0, 0, W, H, fill=th['accent'])
    _tb(sl, sd.get('title','Thank You'), LM, 2.0, CW, 2.2,
        size=52, bold=True, color=th['white'],
        align=PP_ALIGN.CENTER, tracking=-80, lh=1.05)
    _rect(sl, W/2-1.0, 4.35, 2.0, 0.018, fill=th['white'])
    sub = sd.get('subtitle','')
    if sub:
        _tb(sl, sub, LM, 4.65, CW, 0.6,
            size=14, color=th['white'], align=PP_ALIGN.CENTER, lh=1.5)
    if author:
        _tb(sl, author.upper(), LM, 6.5, CW, 0.4,
            size=10, bold=True, color=th['white'],
            align=PP_ALIGN.CENTER, tracking=120)

# =============================================
#  Master Builder
# =============================================
def _build(slides_data, ahex, author, client=None, gen_img=True, prog=None):
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    th = _theme(ahex)
    total = len(slides_data)

    for i, sd in enumerate(slides_data):
        sl = prs.slides.add_slide(blank)
        lay = sd.get('layout', 'split_image')

        if prog:
            prog(i, total, f'Slide {i+1}/{total}: {lay}')

        chart = _make_chart(sd.get('chart_data'), ahex)

        img = None
        ip = sd.get('image_prompt')
        if ip and gen_img and client:
            img = _gen_image(client, ip)
        if img is None and lay in ('title_hero', 'split_image'):
            img = _abstract(ahex)

        if lay == 'title_hero':
            _lay_title(sl, sd, th, ahex, img)
        elif lay == 'stat_callout':
            _lay_stat(sl, sd, th, ahex, chart)
        elif lay == 'split_image':
            _lay_split(sl, sd, th, ahex, img)
        elif lay == 'process_flow':
            _lay_flow(sl, sd, th, ahex)
        elif lay == 'comparison':
            _lay_compare(sl, sd, th, ahex)
        elif lay == 'chart_full':
            _lay_chart(sl, sd, th, ahex, chart)
        elif lay == 'grid_cards':
            _lay_grid(sl, sd, th, ahex)
        elif lay == 'closing':
            _lay_close(sl, sd, th, ahex, author)
        else:
            _lay_split(sl, sd, th, ahex, img)

    buf = io.BytesIO()
    prs.save(buf); buf.seek(0)
    return buf

# =============================================
#  UI
# =============================================
st.markdown('<div class="pg-eyebrow">Slide Generator</div>', unsafe_allow_html=True)
st.markdown(
    '<div class="pg-title">テキストを貼るだけで<br>'
    'コンペを制すスライドが完成する。</div>', unsafe_allow_html=True)
st.markdown(
    '<div class="pg-sub">'
    'AI Content Analysis / Auto Charts / Smart Layout / AI Imagery</div>',
    unsafe_allow_html=True)

# Step 01
st.markdown('<span class="step-label">Step 01 — Gemini API Key</span>',
            unsafe_allow_html=True)
api_key = st.text_input(
    'API Key', type='password',
    placeholder='Gemini API Key (Google AI Studio)',
    label_visibility='collapsed',
    help='https://aistudio.google.com/apikey')

# Step 02
st.markdown('<span class="step-label">Step 02 — Accent Color</span>',
            unsafe_allow_html=True)
accent_label = st.radio('Accent', list(ACCENT_PRESETS.keys()),
                         label_visibility='collapsed')
accent_hex = ACCENT_PRESETS[accent_label]

col1, col2 = st.columns([1, 2])
with col1:
    use_custom = st.checkbox('Custom Color')
with col2:
    if use_custom:
        cx = st.text_input('HEX', value='#C8A96E', label_visibility='collapsed')
        try:
            _rgb(cx); accent_hex = cx
        except:
            st.warning('HEX format error')

st.markdown(
    f'<div style="margin-top:4px;">'
    f'<span class="swatch" style="background:{accent_hex};"></span>'
    f'<span style="font-size:12px;color:#999;letter-spacing:.04em;">{accent_hex}</span>'
    f'</div>', unsafe_allow_html=True)

# Step 03
st.markdown('<span class="step-label">Step 03 — Author / Company</span>',
            unsafe_allow_html=True)
c1, _ = st.columns([1, 1])
with c1:
    author = st.text_input('Author', placeholder='e.g. Creative Division',
                            label_visibility='visible')

# Step 04
st.markdown('<hr class="divider">', unsafe_allow_html=True)
st.markdown('<span class="step-label">Step 04 — Content</span>',
            unsafe_allow_html=True)

text_input = st.text_area(
    'Content', height=360,
    placeholder=(
        '# プレゼンタイトル\n'
        'サブタイトル\n\n'
        '## 市場概況\n'
        '- 市場規模は2024年時点で3.2兆円\n'
        '- 前年比130%の成長率\n\n'
        '## 実現ステップ\n'
        '1. 企画\n2. 設計\n3. 実装\n\n'
        'Gemini / NotebookLM の出力をそのまま貼り付けOK'),
    label_visibility='collapsed')

# Options
st.markdown('<span class="step-label">Options</span>', unsafe_allow_html=True)
gen_images = st.checkbox('AI Image Generation', value=True,
    help='Generate conceptual visuals per slide via Gemini')

st.markdown('<div style="height:6px;"></div>', unsafe_allow_html=True)

# Generate
if st.button('Generate Slides'):
    if not text_input.strip():
        st.warning('Please enter text content.')
    elif not api_key.strip():
        st.warning('Please enter Gemini API Key.')
    else:
        client = _init_gemini(api_key.strip())
        if not client:
            st.error('Failed to initialize Gemini. Check your API key.')
        else:
            bar = st.progress(0, text='Analyzing content...')
            bar.progress(10, text='AI analyzing content structure...')
            slides = _analyze(client, text_input)

            if not slides:
                st.error('Content analysis failed.')
            else:
                bar.progress(25, text=f'{len(slides)} slides structured')

                with st.expander(f'Slide Structure ({len(slides)} slides)',
                                 expanded=False):
                    for i, sd in enumerate(slides):
                        lay = sd.get('layout', '')
                        ttl = sd.get('title', '')
                        st.markdown(
                            f'<div class="preview-card">'
                            f'<div class="preview-tag">{i+1:02d} / {lay}</div>'
                            f'<div class="preview-name">{ttl}</div>'
                            f'</div>', unsafe_allow_html=True)

                def update(i, total, msg):
                    pct = 25 + int(65 * (i+1) / total)
                    bar.progress(pct, text=msg)

                pptx_buf = _build(
                    slides, accent_hex, author,
                    client=client, gen_img=gen_images,
                    prog=update)

                bar.progress(100, text='Complete')

                nt = len(slides)
                layouts = len(set(sd.get('layout','') for sd in slides))
                fname = f'slides_{datetime.now().strftime("%Y%m%d_%H%M")}.pptx'

                st.success(f'{nt} slides / {layouts} layout types / {accent_hex}')
                st.download_button(
                    'Download PowerPoint',
                    data=pptx_buf, file_name=fname,
                    mime='application/vnd.openxmlformats-officedocument.presentationml.presentation')

st.markdown(
    '<div style="text-align:center;font-size:10px;color:#E0E0E0;'
    'letter-spacing:.18em;text-transform:uppercase;margin-top:64px;">'
    'Powered by Gemini + Claude</div>', unsafe_allow_html=True)
"""
Slide Generator — 世界最高峰の広告代理店クオリティ
DTP原則を徹底実装: トラッキング・ベースライングリッド・タイプスケール・行間
"""

import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import io, re
from datetime import datetime

# ─────────────────────────────────────────
st.set_page_config(page_title="Slide Generator", layout="centered",
                   initial_sidebar_state="collapsed")

st.markdown("""
<style>
#MainMenu, footer, header { visibility: hidden; }
.stApp { background:#FFF; font-family:'Helvetica Neue',Arial,sans-serif; }
.main .block-container { max-width:740px; padding-top:56px; padding-bottom:88px; }

.pg-eyebrow { font-size:10px; font-weight:700; letter-spacing:.2em;
              text-transform:uppercase; color:#CCC; margin-bottom:8px; }
.pg-title   { font-size:26px; font-weight:700; letter-spacing:-.3px;
              color:#111; margin-bottom:6px; line-height:1.2; }
.pg-sub     { font-size:13px; color:#999; margin-bottom:52px; line-height:1.7; }
.step-label { font-size:10px; font-weight:700; letter-spacing:.2em;
              text-transform:uppercase; color:#CCC; margin:40px 0 16px; display:block; }

/* Radio cards */
div[data-testid="stRadio"] label { display:none!important; }
div[data-testid="stRadio"] > div { display:flex; flex-direction:column; gap:8px; }
div[data-testid="stRadio"] > div > label {
    border:1px solid #EAEAEA!important; border-radius:2px!important;
    padding:18px 22px!important; background:#FAFAFA!important;
    cursor:pointer!important; transition:all .15s!important;
}
div[data-testid="stRadio"] > div > label:hover {
    border-color:#111!important; background:#FFF!important; }
div[data-testid="stRadio"] > div > label[data-checked="true"] {
    border-color:#111!important; background:#FFF!important;
    box-shadow:inset 3px 0 0 #111!important; }

/* Inputs */
label, .stTextArea label, .stTextInput label {
    font-size:10px!important; font-weight:700!important;
    letter-spacing:.15em!important; text-transform:uppercase!important; color:#BBBBBB!important; }
.stTextArea textarea {
    background:#F9F9F9!important; border:1px solid #E8E8E8!important;
    border-radius:2px!important; font-size:13.5px!important; color:#111!important;
    font-family:'Helvetica Neue',Arial,sans-serif!important;
    line-height:1.7!important; padding:18px!important; resize:vertical!important; }
.stTextArea textarea:focus { border-color:#111!important; box-shadow:none!important; }
.stTextInput input {
    background:#F9F9F9!important; border:1px solid #E8E8E8!important;
    border-radius:2px!important; font-size:13.5px!important; color:#111!important;
    padding:11px 15px!important; }
.stTextInput input:focus { border-color:#111!important; box-shadow:none!important; }

/* Buttons */
.stButton>button {
    background:#111!important; color:#FFF!important; border:none!important;
    border-radius:2px!important; padding:15px 32px!important;
    font-size:11px!important; font-weight:700!important; letter-spacing:.14em!important;
    text-transform:uppercase!important; width:100%!important; transition:opacity .15s!important; }
.stButton>button:hover { opacity:.65!important; }
.stDownloadButton>button {
    background:#FFF!important; color:#111!important;
    border:1.5px solid #111!important; border-radius:2px!important;
    padding:15px 32px!important; font-size:11px!important; font-weight:700!important;
    letter-spacing:.14em!important; text-transform:uppercase!important;
    width:100%!important; margin-top:10px!important; transition:background .15s!important; }
.stDownloadButton>button:hover { background:#F5F5F5!important; }

/* Swatch */
.swatch { display:inline-block; width:13px; height:13px; border-radius:1px;
          margin-right:8px; vertical-align:middle; border:1px solid rgba(0,0,0,.08); }
.divider { border:none; border-top:1px solid #F0F0F0; margin:44px 0; }
[data-testid="column"] { padding:0 6px!important; }
</style>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────
#  ① DTPコアユーティリティ
# ─────────────────────────────────────────
W, H = 13.33, 7.5   # slide size in inches

def rgb(h: str) -> RGBColor:
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def run_style(run, font='Calibri', size=None, bold=False, italic=False,
              color: RGBColor=None, tracking: int=0, ea='Meiryo UI'):
    """
    tracking: hundredths-of-a-point (OOXML `spc`)
      +100 = +1pt (wide, for caps labels)
      -80  = -0.8pt (tight, for large display type)
    """
    run.font.name = font
    if size is not None: run.font.size = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    if color: run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    # Tracking
    if tracking != 0:
        rPr.set('spc', str(int(tracking)))
    # East-Asian font (Japanese support)
    for ch in list(rPr):
        if ch.tag == qn('a:ea'): rPr.remove(ch)
    ea_el = etree.SubElement(rPr, qn('a:ea'))
    ea_el.set('typeface', ea)

def para_style(p, align=PP_ALIGN.LEFT,
               sp_before: float=0, sp_after: float=0,
               line_pct: float=1.0):
    """
    line_pct: 1.0=single, 1.4=140%...
    """
    p.alignment  = align
    p.space_before = Pt(sp_before)
    p.space_after  = Pt(sp_after)
    if line_pct != 1.0:
        pPr = p._p.get_or_add_pPr()
        for ch in list(pPr):
            if ch.tag == qn('a:lnSpc'): pPr.remove(ch)
        lnSpc  = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_pct * 100000)))

def shape_rect(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    s.line.fill.background()
    return s

def textbox(slide, text, l, t, w, h,
            font='Calibri', size=16, bold=False, italic=False,
            color: RGBColor=None, align=PP_ALIGN.LEFT,
            tracking: int=0, line_pct: float=1.0,
            sp_before: float=0, wrap=True, ea='Meiryo UI'):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    para_style(p, align=align, sp_before=sp_before, line_pct=line_pct)
    r = p.add_run(); r.text = text
    run_style(r, font=font, size=size, bold=bold, italic=italic,
              color=color, tracking=tracking, ea=ea)
    return tb

# ─────────────────────────────────────────
#  ② タイプスケール (黄金比ベース)
# ─────────────────────────────────────────
# Display  → 60pt  tracking -100
# H1       → 56pt  tracking -80
# H2       → 32pt  tracking -40
# Body     → 16pt  tracking   0   line-height 1.45
# Sub      → 13pt  tracking   0   line-height 1.38
# Caption  → 10pt  tracking +100
TS = dict(
    display  = dict(size=60, bold=True,  tracking=-100, line_pct=1.05),
    h1       = dict(size=56, bold=True,  tracking=-80,  line_pct=1.05),
    h2       = dict(size=32, bold=True,  tracking=-40,  line_pct=1.1),
    h2_light = dict(size=30, bold=False, tracking=-20,  line_pct=1.1),
    body     = dict(size=16, bold=False, tracking=0,    line_pct=1.45),
    sub      = dict(size=13, bold=False, tracking=0,    line_pct=1.38),
    caption  = dict(size=10, bold=False, tracking=100,  line_pct=1.0),
    label    = dict(size=10, bold=True,  tracking=120,  line_pct=1.0),
    num      = dict(size=10, bold=True,  tracking=80,   line_pct=1.0),
)

# Margin standard
LM = 0.75   # left margin
RM = 0.75   # right margin
CW = W - LM - RM   # content width = 11.83

# ─────────────────────────────────────────
#  ③ コンテンツブロック（全テーマ共通）
# ─────────────────────────────────────────
def content_block(slide, items, l, t, w, h, th):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True

    for i, item in enumerate(items):
        p  = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        lv = item.get('level', 0)
        ib = item.get('bold', False)
        txt = item['text']

        if lv == 0 and ib:
            # ── アクセントヘッダー (### 見出し)
            para_style(p, sp_before=18 if i > 0 else 4,
                       sp_after=2, line_pct=1.25)
            r = p.add_run(); r.text = txt
            run_style(r, font='Calibri', size=17, bold=True,
                      color=th['accent'], tracking=-20, ea='Meiryo UI')

        elif lv == 0:
            # ── 主要箇条書き
            para_style(p, sp_before=9 if i > 0 else 0,
                       sp_after=0, line_pct=1.45)
            r = p.add_run()
            r.text = '\u2014\u2002' + txt   # —  + thin space
            run_style(r, font='Calibri', size=16, bold=False,
                      color=th['body'], tracking=0, ea='Meiryo UI')

        else:
            # ── サブ箇条書き (インデント)
            para_style(p, sp_before=4, sp_after=0, line_pct=1.38)
            r = p.add_run()
            r.text = '\u2003\u00B7\u2002' + txt  # em-space + middle-dot + thin-space
            run_style(r, font='Calibri', size=13, bold=False,
                      color=th['sub'], tracking=0, ea='Meiryo UI')

# ─────────────────────────────────────────
#  ④ テーマ設定
# ─────────────────────────────────────────
ACCENT_PRESETS = {
    'ネイビー':   '#0A2463',
    'ゴールド':   '#A67C3A',
    'クリムゾン': '#8C1C2E',
    'エメラルド': '#0B5E4A',
    'スレート':   '#2E3F4F',
}

def get_theme(key: str, accent_hex: str) -> dict:
    a = rgb(accent_hex)
    base = dict(
        accent = a,
        white  = rgb('FFFFFF'),
        ink    = rgb('0D0D0D'),    # near-black (softer than pure black)
        body   = rgb('2C2C2C'),    # body text
        sub    = rgb('717171'),    # sub-bullets / captions
        rules  = rgb('EBEBEB'),    # hairlines
        num_c  = rgb('CCCCCC'),    # slide numbers
        muted  = rgb('9A9A9A'),    # subtitles / dates
    )
    if key == 'MINIMAL':
        base.update(title_bg=rgb('FFFFFF'), hdr_bg=rgb('FFFFFF'), hdr_txt=rgb('0D0D0D'))
    elif key == 'NOIR':
        base.update(title_bg=rgb('0C0C0C'), hdr_bg=rgb('FFFFFF'), hdr_txt=rgb('0D0D0D'))
    elif key == 'BOLD':
        base.update(title_bg=rgb('FFFFFF'), hdr_bg=a,             hdr_txt=rgb('FFFFFF'))
    return base

# ─────────────────────────────────────────
#  ⑤ MINIMAL — 知性のグリッドシステム
#    大型ナンバリング × 水平リズム × 広大な余白
# ─────────────────────────────────────────
def minimal_title(sl, sd, th, author, today):
    # ─ 最上部ルール (アクセント)
    shape_rect(sl, 0, 0, W, 0.045, fill=th['accent'])

    # ─ 著者・日付 (キャプション, wide tracking)
    if author:
        textbox(sl, author.upper(), LM, 0.30, 5.0, 0.35,
                size=10, bold=True, color=th['muted'], tracking=120)
    textbox(sl, today, W-RM-3.0, 0.30, 3.0, 0.35,
            size=10, color=th['muted'], align=PP_ALIGN.RIGHT, tracking=80)

    # ─ タイトル (display, tight tracking)
    textbox(sl, sd['title'], LM, 2.05, CW, 1.75,
            size=60, bold=True, color=th['ink'], tracking=-100, line_pct=1.05)

    # ─ アクセントライン (タイトル直後)
    shape_rect(sl, LM, 3.95, 2.4, 0.022, fill=th['accent'])

    # ─ サブタイトル
    if sd.get('subtitle'):
        textbox(sl, sd['subtitle'], LM, 4.18, CW*0.75, 0.65,
                size=17, color=th['muted'], tracking=0, line_pct=1.5)

    # ─ 最下部ルール
    shape_rect(sl, 0, H-0.045, W, 0.045, fill=th['accent'])


def minimal_content(sl, sd, th, num):
    # ─ スライド番号 (小さく、右上、refined)
    textbox(sl, f'{num:02d}', W-RM-0.8, 0.22, 0.8, 0.38,
            size=10, bold=True, color=th['num_c'],
            align=PP_ALIGN.RIGHT, tracking=80)

    # ─ タイトル (H2 bold)
    textbox(sl, sd.get('title',''), LM, 0.20, CW-1.0, 0.82,
            size=32, bold=True, color=th['ink'], tracking=-40, line_pct=1.1)

    # ─ ヘアラインルール (全幅 / アクセント部分 2.2")
    shape_rect(sl, LM, 1.15, CW, 0.012, fill=th['rules'])
    shape_rect(sl, LM, 1.15, 2.2, 0.012, fill=th['accent'])

    # ─ コンテンツ
    items = sd.get('content', [])
    if items:
        content_block(sl, items, LM, 1.35, CW, H-1.35-0.45, th)

    # ─ 最下部ルール
    shape_rect(sl, 0, H-0.045, W, 0.045, fill=th['accent'])


# ─────────────────────────────────────────
#  ⑥ NOIR — シネマティック・ラグジュアリー
#    漆黒タイトル × 白コンテンツ × 鯯利なアクセント
# ─────────────────────────────────────────
def noir_title(sl, sd, th, author, today):
    # ─ 全面黒背景
    shape_rect(sl, 0, 0, W, H, fill=th['title_bg'])

    # ─ トップ: 極細アクセントライン
    shape_rect(sl, 0, 0, W, 0.045, fill=th['accent'])

    # ─ 右側: 装飾的ヴァーティカルライン (非常に細い、ほぼ見えない程度)
    shape_rect(sl, 12.15, 0.4, 0.018, H-0.8, fill=rgb('1C1C1C'))

    # ─ タイトル (display, 白, tight tracking)
    textbox(sl, sd['title'], LM, 1.70, 10.5, 2.0,
            size=60, bold=True, color=th['white'],
            tracking=-100, line_pct=1.05)

    # ─ アクセントライン (タイトル直後)
    shape_rect(sl, LM, 3.82, 2.0, 0.022, fill=th['accent'])

    # ─ サブタイトル
    if sd.get('subtitle'):
        textbox(sl, sd['subtitle'], LM, 4.06, 9.5, 0.65,
                size=17, color=th['muted'], tracking=0, line_pct=1.5)

    # ─ フッター
    parts = [p for p in [author, today] if p]
    if parts:
        textbox(sl, '  /  '.join(parts), LM, 7.02, CW, 0.28,
                size=10, color=rgb('3A3A3A'),
                align=PP_ALIGN.RIGHT, tracking=80)

    # ─ ボトムライン
    shape_rect(sl, 0, H-0.045, W, 0.045, fill=th['accent'])


def noir_content(sl, sd, th, num):
    # ─ 左: 極細アクセントストリップ
    shape_rect(sl, 0, 0, 0.048, H, fill=th['accent'])

    # ─ スライド番号 (右上, muted)
    textbox(sl, f'{num:02d}', W-RM-0.85, 0.20, 0.85, 0.40,
            size=10, bold=True, color=th['num_c'],
            align=PP_ALIGN.RIGHT, tracking=80)

    # ─ タイトル (H2 bold)
    textbox(sl, sd.get('title',''), LM, 0.20, CW-1.0, 0.82,
            size=32, bold=True, color=th['ink'], tracking=-40, line_pct=1.1)

    # ─ ダブルヘアライン (全幅グレー + 短めアクセント)
    shape_rect(sl, LM, 1.15, CW, 0.012, fill=th['rules'])
    shape_rect(sl, LM, 1.15, 1.8, 0.012, fill=th['accent'])

    # ─ コンテンツ
    items = sd.get('content', [])
    if items:
        content_block(sl, items, LM, 1.35, CW, H-1.35-0.45, th)

    # ─ ボトムライン
    shape_rect(sl, 0, H-0.045, W, 0.045, fill=th['accent'])


# ─────────────────────────────────────────
#  ⑦ BOLD — パネル分割 × ブランド刈印
#    クリエイティブエージェンシーの大胆な構成力
# ─────────────────────────────────────────
PANEL_W = 5.1   # left panel width

def bold_title(sl, sd, th, author, today):
    # ─ 左パネル (アクセントカラー)
    shape_rect(sl, 0, 0, PANEL_W, H, fill=th['accent'])

    # ─ パネル内: 著者 (最上部, 白, wide tracking)
    if author:
        textbox(sl, author.upper(), 0.55, 0.30, PANEL_W-0.65, 0.38,
                size=9, bold=True, color=rgb('FFFFFF'), tracking=120)
    else:
        textbox(sl, today, 0.55, 0.30, PANEL_W-0.65, 0.38,
                size=9, color=rgb('FFFFFF'), tracking=80)

    # ─ パネル内: アクセントライン (水平)
    shape_rect(sl, 0.55, 1.88, 1.8, 0.022, fill=rgb('FFFFFF'))

    # ─ パネル内: タイトル (白, large)
    textbox(sl, sd['title'], 0.55, 2.05, PANEL_W-0.7, 3.5,
            size=44, bold=True, color=th['white'],
            tracking=-80, line_pct=1.08)

    # ─ 右パネル: サブタイトル (中央寄り縦位置)
    if sd.get('subtitle'):
        textbox(sl, sd['subtitle'], PANEL_W+0.65, 2.2, W-PANEL_W-1.2, 2.0,
                size=20, color=th['ink'], tracking=-10, line_pct=1.5)

    # ─ 右パネル: 日付
    textbox(sl, today if author else '', W-RM-2.5, 7.0, 2.5, 0.30,
            size=10, color=th['muted'], align=PP_ALIGN.RIGHT, tracking=80)


def bold_content(sl, sd, th, num):
    # ─ ヘッダーバー (アクセントカラー)
    shape_rect(sl, 0, 0, W, 1.05, fill=th['hdr_bg'])

    # ─ タイトル (ヘッダー内, 白)
    textbox(sl, sd.get('title',''), LM, 0.18, CW-1.2, 0.72,
            size=26, bold=True, color=th['hdr_txt'],
            tracking=-30, line_pct=1.1)

    # ─ スライド番号 (ヘッダー右端, 白, 大きめ)
    textbox(sl, f'{num:02d}', W-RM-0.9, 0.24, 0.9, 0.5,
            size=18, bold=True,
            color=rgb('FFFFFF') if th['hdr_txt'] == th['white'] else th['num_c'],
            align=PP_ALIGN.RIGHT, tracking=60)

    # ─ コンテンツ
    items = sd.get('content', [])
    if items:
        content_block(sl, items, LM, 1.22, CW, H-1.22-0.45, th)

    # ─ ボトムライン
    shape_rect(sl, 0, H-0.045, W, 0.045, fill=th['accent'])


# ─────────────────────────────────────────
#  ⑧ マスタービルダー
# ─────────────────────────────────────────
TITLE_FN   = {'MINIMAL': minimal_title, 'NOIR': noir_title, 'BOLD': bold_title}
CONTENT_FN = {'MINIMAL': minimal_content,'NOIR': noir_content,'BOLD': bold_content}

def build_pptx(slides_data, theme_key, accent_hex, author='') -> io.BytesIO:
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    th    = get_theme(theme_key, accent_hex)
    today = datetime.now().strftime('%Y.%m.%d')
    cnum  = 0

    for sd in slides_data:
        sl = prs.slides.add_slide(blank)
        if sd['type'] == 'title':
            TITLE_FN[theme_key](sl, sd, th, author, today)
        else:
            cnum += 1
            CONTENT_FN[theme_key](sl, sd, th, cnum)

    buf = io.BytesIO()
    prs.save(buf); buf.seek(0)
    return buf

# ─────────────────────────────────────────
#  ⑨ テキストパーサー
# ─────────────────────────────────────────
def parse_slides(text: str) -> list[dict]:
    lines = text.strip().split('\n')
    slides, cur = [], [None]

    def push():
        c = cur[0]
        if c and (c.get('title') or c.get('content')): slides.append(c)
        cur[0] = None

    def cs(title=''):
        return {'type':'content','title':title,'content':[]}

    for raw in lines:
        s = raw.strip()
        if re.match(r'^(-{3,}|={3,})$', s):            push(); continue
        if not s:                                        continue
        if s.startswith('# '):
            push(); t = s[2:].strip()
            cur[0] = ({'type':'title','title':t,'subtitle':'','content':[]}
                      if not slides else cs(t)); continue
        if s.startswith('## '):
            push(); cur[0] = cs(s[3:].strip()); continue
        if s.startswith('### '):
            if not cur[0]: cur[0] = cs(s[4:].strip())
            else: cur[0]['content'].append({'level':0,'text':s[4:].strip(),'bold':True})
            continue
        m = re.match(r'^[-*・•]\s+(.+)', s)
        if m:
            if not cur[0]: cur[0] = cs()
            cur[0]['content'].append({'level':0,'text':m.group(1).strip()}); continue
        m2 = re.match(r'^\s{2,}[-*・•]\s+(.+)', raw)
        if m2:
            if not cur[0]: cur[0] = cs()
            cur[0]['content'].append({'level':1,'text':m2.group(1).strip()}); continue
        m3 = re.match(r'^\d+[.。)]\s+(.+)', s)
        if m3:
            if not cur[0]: cur[0] = cs()
            cur[0]['content'].append({'level':0,'text':m3.group(1).strip()}); continue
        c = cur[0]
        if not c:   cur[0] = {'type':'title','title':s,'subtitle':'','content':[]}
        elif c['type'] == 'title' and not c.get('subtitle'): c['subtitle'] = s
        else: c['content'].append({'level':0,'text':s})

    push()
    return slides or [{'type':'title','title':'Presentation','subtitle':'','content':[]}]

# ─────────────────────────────────────────
#  ⑩ UI
# ─────────────────────────────────────────
st.markdown('<div class="pg-eyebrow">Slide Generator</div>', unsafe_allow_html=True)
st.markdown('<div class="pg-title">テキストを貿るだけで<br>コンペを制すスライドが完成する。</div>',
            unsafe_allow_html=True)
st.markdown('<div class="pg-sub">DTP原則を徹底実装。トラッキング・ベースライングリッド・タイプスケール・行間比率を自動制御。</div>',
            unsafe_allow_html=True)

# ── Step 01: テーマ ───────────────────────
st.markdown('<span class="step-label">Step 01 — Design Language</span>', unsafe_allow_html=True)

theme_raw = st.radio('テーマ', [
    'MINIMAL  ——  グリッドとナンバリング。知性を纚ったコーポレートデザイン。',
    'NOIR  ——  漆黒のタイトルスライド。圧倒的な存在感でコンペを制す。',
    'BOLD  ——  パネル分割と大型タイポグラフィー。ブランドを刈む。',
], label_visibility='collapsed')
theme_key = theme_raw.split()[0]

# ── Step 02: アクセントカラー ─────────────────
st.markdown('<span class="step-label">Step 02 — Accent Color</span>', unsafe_allow_html=True)

accent_label = st.radio('アクセントカラー',
    list(ACCENT_PRESETS.keys()),
    label_visibility='collapsed')
accent_hex = ACCENT_PRESETS[accent_label]

col1, col2 = st.columns([1, 2])
with col1:
    use_custom = st.checkbox('カスタムカラー')
with col2:
    if use_custom:
        cx = st.text_input('HEX', value='#C8A96E', label_visibility='collapsed')
        try: rgb(cx); accent_hex = cx
        except: st.warning('HEX形式エラー（例: #C8A96E）')

st.markdown(
    f'<div style="margin-top:4px;">'
    f'<span class="swatch" style="background:{accent_hex};"></span>'
    f'<span style="font-size:12px;color:#999;letter-spacing:.04em;">{accent_hex}</span>'
    f'</div>', unsafe_allow_html=True)

# ── Step 03: 作成者 ───────────────────────
st.markdown('<span class="step-label">Step 03 — Author / Company</span>', unsafe_allow_html=True)
c1, _ = st.columns([1, 1])
with c1:
    author = st.text_input('組織名・作成者',
                           placeholder='例：Creative Division',
                           label_visibility='visible')

# ── Step 04: テキスト入力 ────────────────
st.markdown('<hr class="divider">', unsafe_allow_html=True)
st.markdown('<span class="step-label">Step 04 — Content</span>', unsafe_allow_html=True)

text_input = st.text_area(
    'テキスト入力', height=360,
    placeholder=(
        '# プレゼンテーションタイトル\n'
        'サブタイトル・キャッチコピー\n\n'
        '## スライドの見出し\n'
        '- 伝えたいメッセージ\n'
        '- 重要なポイント\n'
        '  - 詳細・補足\n\n'
        '## 次のスライド\n'
        '### 強調したい小見出し\n'
        '- 項目A\n'
        '- 項目B\n\n'
        '---  （スライドを手動で区切る）'
    ), label_visibility='collapsed')

with st.expander('テキストの書き方ガイド'):
    st.markdown("""
| 書き方 | 効果 |
|--------|------|
| `# タイトル` | タイトルスライド |
| `## 見出し` | コンテンツスライド見出し |
| `### 小見出し` | アクセントカラーで強調 |
| `- 項目` | 箇条書き (—  マーカー) |
| `  - サブ` | インデント箇条書き (·  マーカー) |
| `---` | スライド強制分割 |

GeminiやNotebookLMの出力をそのまま貼り付けてOKです。
""")

st.markdown('<div style="height:6px;"></div>', unsafe_allow_html=True)

if st.button('スライドを生成する'):
    if not text_input.strip():
        st.warning('テキストを入力してください。')
    else:
        with st.spinner('生成中...'):
            slides   = parse_slides(text_input)
            pptx_buf = build_pptx(slides, theme_key, accent_hex, author)

        nt = len(slides)
        nc = sum(1 for s in slides if s['type'] != 'title')
        fname = f'slides_{datetime.now().strftime("%Y%m%d_%H%M")}.pptx'

        st.success(
            f'生成完了 — {nt} 枚（コンテンツ {nc} 枚）  '
            f'[ {theme_key}  ×  {accent_hex} ]'
        )
        st.download_button(
            'PowerPoint をダウンロード', data=pptx_buf, file_name=fname,
            mime='application/vnd.openxmlformats-officedocument.presentationml.presentation')

st.markdown(
    '<div style="text-align:center;font-size:10px;color:#E0E0E0;'
    'letter-spacing:.18em;text-transform:uppercase;margin-top:64px;">'
    'Powered by Claude</div>', unsafe_allow_html=True)
